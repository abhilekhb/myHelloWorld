from datetime import datetime 
from pptx import Presentation
from pptx.util import Inches
#import argparse - unused so far 
import os 
import math
import pandas as pd
import numpy as np
from datetime import date
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor, ColorFormat
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path
import win32com.client
from pptx.util import Pt

#"""
def analyze_ppt(input, output):
    # Take the input file and analyze the structure.
    #The output file contains marked up information to make it easier
    #for generating future powerpoint templates.

    prs = Presentation(input)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if 'Title' not in shape.text:
                        shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                #print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)

def df_to_table(slide, df, left, top, width, height, colnames=None):
    """Converts a Pandas DataFrame to a standard PowerPoint table for a PPT file 
    Arguments:
     - slide: slide object from pptx lib which has the slide on which the table should appear
     - df: Pandas DataFrame with the data
    Optional arguments:
     - colnames
     """
    rows, cols = df.shape
    #print('rows=',rows,'cols=',cols)
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = " ".join(col_name)
        res.table.cell(0, col_index).text = col_name
        paragraph = res.table.cell(0, col_index).text_frame.paragraphs[0]
        paragraph.font.size = Pt(15)
        paragraph.alignment = PP_ALIGN.CENTER
        res.table.cell(0, col_index).fill.solid()
        res.table.cell(0, col_index).fill.fore_color.rgb = RGBColor(255,100,0)
        #print(col_name)

    m = df.to_numpy()
    #print('m numpy array:',m)


    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text
            paragraph = res.table.cell(row+1, col).text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(0, 0, 0) # use black color for now 
            res.table.cell(row+1, col).fill.background()

template_ip_ppt = 'C:\\Users\\abhardwaj\\Desktop\\Training\\Python\\Code\\input_ppt.pptx'
template_op_ppt = 'C:\\Users\\abhardwaj\\Desktop\\Training\\Python\\Code\\output_ppt.pptx'
#analyze_ppt(template_ip_ppt,template_op_ppt)

template_ip_ppt = 'C:\\Users\\abhardwaj\\Desktop\\Training\\Python\\Code\\input_ppt.pptx'
final_op_ppt = 'C:\\Users\\abhardwaj\\Desktop\\Training\\Python\\Code\\final_output_ppt.pptx'
srcDir = 'C:\\Users\\abhardwaj\\Desktop\\Dashboards\\Service Catalogues for Digital Products\\'
list_of_files = []
prs = Presentation(template_ip_ppt)

# Create a title slide first
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
pic = slide.shapes.add_picture('C:\\Users\\abhardwaj\\Desktop\\Training\\Python\\Code\\servops.jpg',Inches(1),Inches(1))
pic.rotation = -45.0 # display the picture rotated by -45 (360-45=315) degrees 
#print('no of shapes = ',len(slide.shapes))
#shape1 = slide.shapes[0]
#shape2 = slide.shapes[1]
#print('shape1 =',shape1.name, 'shape id =', shape1.shape_id, 'shape type =',shape1.shape_type)
#print('shape2 =',shape2.name, 'shape id =', shape2.shape_id, 'shape type =',shape2.shape_type)
#if shape1.has_text_frame:
#    print('shape1 has a text frame!')
#if shape2.has_text_frame:
#    print('shape2 has a text frame!')
title = slide.shapes.title
title.text = "Service Ops Dashboards"
para1 = title.text_frame.paragraphs[0]
para1.font.name = 'Calibri'
para1.font.color.rgb = RGBColor(255, 100, 0) # print in Orange color 
#title.line.color.rgb = RGBColor(0x00, 0x16, 0xBC)

subtitle = slide.placeholders[11]
"""
subtitle.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
text_frame = subtitle.text_frame
p = text_frame.paragraphs[0]
font.name = 'Calibri'
"""
subtitle.text = "Auto-Generated on {:%d %B-%Y, %I:%M:%S %p}".format(datetime.now())
para1 = subtitle.text_frame.paragraphs[0]
para1.font.name = 'Calibri'
para1.font.size = Pt(12)
para1.font.italic = True
para1.font.color.rgb = RGBColor(255, 0, 0) # print in red color 

#test code to insert a rectangle shape on the slide 
#shapes = slide.shapes
#left = top = width = height = Inches(1.0)
#shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

#"""
for r, d, f in os.walk(srcDir):
    for file in f:
        fullPath = os.path.join(r, file)
        # skip if it is symbolic link or if it is a temporary windows file starting with '~$' symbol
        if not (os.path.islink(fullPath)):
            if not file.startswith('~$'):
                if file.endswith('xlsx') :
                    list_of_files.append(fullPath)
                    #print(fullPath)
                    #data = pd.read_excel(fullPath)
                    xls = pd.ExcelFile(fullPath)
                    # list all sheets in the file
                    print(xls.sheet_names)
                    # to read just one sheet to dataframe:
                    # Create a section overview slide first 
                    title_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(title_slide_layout)
                    #print('no of shapes = ',len(slide.shapes))
                    #shape1 = slide.shapes[0]
                    #shape2 = slide.shapes[1]
                    #print('shape1 =',shape1.name, 'shape id =', shape1.shape_id)
                    #print('shape2 =',shape2.name, 'shape id =', shape2.shape_id)
                    title = slide.shapes.title
                    para2 = title.text_frame.paragraphs[0]
                    para2.font.size = Pt(12)
                    title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    title.text_frame.word_wrap = False 
                    title.line.color.rgb = RGBColor(0x00, 0x16, 0xBC)
                    title.text = xls.sheet_names[0] + " Services Dashboard"
                    subtitle = slide.placeholders[11]
                    subtitle.text = "Auto-Generated on {:%d %B-%Y, %I:%M:%S %p}".format(datetime.now())
                    para2 = subtitle.text_frame.paragraphs[0]
                    para2.font.size = Pt(12)
                    para2.font.italic = True
                    para2.font.color.rgb = RGBColor(255, 0, 0) # print in red color 
                    data = pd.read_excel(fullPath, sheet_name=xls.sheet_names[0])
                    df1 = pd.DataFrame(data,columns=['Product functionality','Service Ops Availability','Service Ops Performance'])
                    #remove all rows from the dataframe that contains certain strings e.g., 'Not sure ' or 'Duplicate?'
                    df1 = df1[~df1['Service Ops Availability'].isin(['Not sure','Not sure ','Duplicate','Duplicate?'])]
                    #remove all rows from the dataframe where a certain column is null or NaN
                    df1 = df1[pd.notnull(df1['Service Ops Availability'])] 
                    #replace values in the other dataframe columns which contain a null or a NaN
                    values = {'Product functionality':'Not provided', 'Service Ops Performance': 'Not Required'}
                    df1.fillna(value=values,inplace=True)
                    #df1 = df1.dropna()
                    #print(df1)
                    #insert an empty rectangle on the slide 
                    #shapes = slide.shapes
                    #height = Inches(1.2)
                    #top = Inches(3.2)
                    #left = Inches(0.45)
                    #width = Inches(7.75)
                    #shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                    #fill = shape.fill
                    #line = shape.line
                    #line.color.rgb = RGBColor(255, 0, 0) #red color
                    #line.width = Pt(5.5) 
                    #fill.solid() # for solid fill 
                    #fill.fore_color.rgb = RGBColor(255, 0, 0) # for red color fill 
                    #fill.transparency = 0.95 # this doesn't seem to work for now although it doesn't give error 
                    #shape.fill.background() #for empty shape with no fill color 
                    prs.save(final_op_ppt)
                    # Create the detailed slides of this section
                    row_count = len(df1)
                    #print('row_count=',row_count)
                    row_count = math.ceil(row_count/10)
                    #print('row_count=',row_count)
                    times = 1
                    i = 0  # starting with the 1st row in the Data Frame 
                    j = 10 # fetch 9 rows from the Data Frame at a time 
                    while times<=row_count:
                        title_slide_layout = prs.slide_layouts[4]
                        slide = prs.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        subtitle = slide.placeholders[11]
                        #txt_content = slide.placeholders[10]
                        title.text = xls.sheet_names[0] + ' Services Dashboard [' + str(times) + ']'
                        subtitle.text = "Auto-Generated on {:%d %B-%Y, %I:%M:%S %p}".format(datetime.now())
                        #txt_content.text = df1.to_string(columns=[0:1],na_rep="Empty")
                        #txt_content.text = df1[i:j].to_string(na_rep="Empty")
                        top = Inches(1.5)
                        left = Inches(0.5)
                        width = Inches(12)
                        height = Inches(5.0)
                        df_to_table(slide, df1[i:j], left, top, width, height, colnames=None)
                        prs.save(final_op_ppt)
                        times += 1 # prep for the next iteration 
                        i += 10    # fetch the next block of 10 rows starting at i position 
                        j += 10    # fetch the next block of 10 rows ending at j position 

print('Total',len(list_of_files),'files processed into powerpoint')
#announce that the program has been successfully completed 
#speaker = win32com.client.Dispatch("SAPI.SpVoice") 
#speaker.Speak('Program successfully completed. Goodbye!')
	
def ppt2pdf(ppt_target_file):
    file_path = Path(ppt_target_file).resolve()
    out_file = file_path.parent / file_path.stem
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    pdf = powerpoint.Presentations.Open(file_path, WithWindow=False)
    pdf.SaveAs(out_file, 32)
    pdf.Close()
    powerpoint.Quit()

#ppt2pdf(final_op_ppt) # works but a little slower 
#"""
